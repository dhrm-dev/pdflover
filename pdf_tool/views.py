import os
import pandas as pd
import pdfplumber
from django.shortcuts import render
from django.http import HttpResponse, JsonResponse
from django.conf import settings
import tempfile
from PIL import Image
import io
import base64
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from pdf2image import convert_from_path
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from pdf2docx import Converter
import img2pdf
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from django.views.decorators.csrf import csrf_exempt
import json
import pytesseract
import os, uuid
import os, io, tempfile
import camelot
from openpyxl.utils.dataframe import dataframe_to_rows


# Tesseract path set karo (Windows ke liye)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'



# ==============================
# PDF TO EXCEL (MAIN VIEW)
# ==============================
def pdf_to_excel(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name

        try:
            is_computer_generated = check_if_computer_generated(temp_pdf_path)

            if is_computer_generated:
                excel_data = extract_text_from_pdf(temp_pdf_path)
            else:
                excel_data = embed_images_in_excel(temp_pdf_path)

            response = HttpResponse(
                excel_data,
                content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
            response['Content-Disposition'] = 'attachment; filename="converted_excel.xlsx"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)

        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)

    return render(request, 'pdf_to_excel.html')


# ==============================
# CHECK TEXT / SCANNED PDF
# ==============================
def check_if_computer_generated(pdf_path):
    try:
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
            text_pages = 0

            for page in pdf.pages:
                text = page.extract_text()
                if text and len(text.strip()) > 50:
                    text_pages += 1

            return (text_pages / total_pages) > 0.5 if total_pages else False
    except Exception:
        return False


# ==============================
# TEXT PDF ➜ HIGH QUALITY EXCEL
# (REPLACED WITH CAMELOT)
# ==============================
def extract_text_from_pdf(pdf_path):
    wb = Workbook()
    wb.remove(wb.active)

    # 1️⃣ Best quality (border tables)
    tables = camelot.read_pdf(
        pdf_path,
        pages="all",
        flavor="lattice"
    )

    # 2️⃣ Fallback (borderless tables)
    if tables.n == 0:
        tables = camelot.read_pdf(
            pdf_path,
            pages="all",
            flavor="stream"
        )

    if tables.n == 0:
        raise Exception("No tables detected in PDF")

    for i, table in enumerate(tables):
        df = table.df
        sheet_name = f"Table_{i + 1}"[:31]
        ws = wb.create_sheet(title=sheet_name)

        for row in dataframe_to_rows(df, index=False, header=True):
            ws.append(row)

        # Auto column width (polish)
        for col in ws.columns:
            max_len = max(
                len(str(cell.value)) if cell.value else 0
                for cell in col
            )
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 40)

    temp_excel_path = "temp_text_pdf.xlsx"
    wb.save(temp_excel_path)

    with open(temp_excel_path, "rb") as f:
        excel_data = f.read()

    os.unlink(temp_excel_path)
    return excel_data


# ==============================
# SCANNED PDF ➜ IMAGE EXCEL
# (NO OCR)
# ==============================
def embed_images_in_excel(pdf_path):
    wb = Workbook()
    wb.remove(wb.active)

    images = convert_from_path(pdf_path, dpi=200)

    for page_num, image in enumerate(images):
        ws = wb.create_sheet(title=f"Page_{page_num + 1}"[:31])

        img_bytes = io.BytesIO()
        image.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        xl_img = XLImage(img_bytes)
        ws.add_image(xl_img, "A1")

        ws.cell(row=1, column=10, value=f"Scanned PDF – Page {page_num + 1}")

    temp_excel_path = "temp_scanned_pdf.xlsx"
    wb.save(temp_excel_path)

    with open(temp_excel_path, "rb") as f:
        excel_data = f.read()

    os.unlink(temp_excel_path)
    return excel_data


# PDF to Word
import os
import uuid
import tempfile
import pdfplumber
from docx import Document
from django.http import HttpResponse
from django.shortcuts import render

def pdf_to_word(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']

        file_id = str(uuid.uuid4())
        pdf_path = f"/tmp/{file_id}.pdf"
        docx_path = f"/tmp/{file_id}.docx"

        # Save PDF
        with open(pdf_path, 'wb') as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)

        try:
            doc = Document()

            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        doc.add_paragraph(text)

            doc.save(docx_path)

            with open(docx_path, 'rb') as f:
                response = HttpResponse(
                    f.read(),
                    content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                )
                response['Content-Disposition'] = 'attachment; filename="converted_document.docx"'
                return response

        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)

        finally:
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            if os.path.exists(docx_path):
                os.remove(docx_path)

    return render(request, 'pdf_to_word.html')


# PDF to PowerPoint
def pdf_to_ppt(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            # Convert PDF to images
            images = convert_from_path(temp_pdf_path, dpi=150)
            prs = Presentation()
            
            for i, image in enumerate(images):
                # Add slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Convert image to bytes
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format='PNG')
                img_byte_arr.seek(0)
                
                # Add image to slide
                left = Inches(1)
                top = Inches(1)
                pic = slide.shapes.add_picture(img_byte_arr, left, top, height=Inches(5.5))
                
                # Add page number
                txBox = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(2), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"Page {i+1}"
            
            temp_ppt_path = 'temp_converted.pptx'
            prs.save(temp_ppt_path)
            
            with open(temp_ppt_path, 'rb') as f:
                ppt_data = f.read()
            
            response = HttpResponse(
                ppt_data,
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = 'attachment; filename="converted_presentation.pptx"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_converted.pptx'):
                os.unlink('temp_converted.pptx')
    
    return render(request, 'pdf_to_ppt.html')

# PDF to Images
def pdf_to_images(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            images = convert_from_path(temp_pdf_path, dpi=200)
            
            # Create zip file with all images
            import zipfile
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
                for i, image in enumerate(images):
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='PNG')
                    zip_file.writestr(f'page_{i+1}.png', img_byte_arr.getvalue())
            
            zip_buffer.seek(0)
            response = HttpResponse(
                zip_buffer.getvalue(),
                content_type='application/zip'
            )
            response['Content-Disposition'] = 'attachment; filename="converted_images.zip"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
    
    return render(request, 'pdf_to_images.html')

# Merge PDF
def merge_pdf(request):
    if request.method == 'POST' and request.FILES.getlist('pdf_files'):
        pdf_files = request.FILES.getlist('pdf_files')
        
        if len(pdf_files) < 2:
            return HttpResponse("Please select at least 2 PDF files to merge", status=400)
        
        merger = PdfMerger()
        temp_files = []
        
        try:
            for pdf_file in pdf_files:
                temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                for chunk in pdf_file.chunks():
                    temp_pdf.write(chunk)
                temp_pdf.close()
                temp_files.append(temp_pdf.name)
                merger.append(temp_pdf.name)
            
            merged_pdf_path = 'temp_merged.pdf'
            merger.write(merged_pdf_path)
            merger.close()
            
            with open(merged_pdf_path, 'rb') as f:
                merged_data = f.read()
            
            response = HttpResponse(
                merged_data,
                content_type='application/pdf'
            )
            response['Content-Disposition'] = 'attachment; filename="merged_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            for temp_file in temp_files:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            if os.path.exists('temp_merged.pdf'):
                os.unlink('temp_merged.pdf')
    
    return render(request, 'merge_pdf.html')

# Split PDF
def split_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        pages = request.POST.get('pages', '')  # e.g., "1,3,5-8"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            total_pages = len(reader.pages)
            
            # Parse pages input
            selected_pages = parse_page_selection(pages, total_pages)
            
            if not selected_pages:
                return HttpResponse("Invalid page selection", status=400)
            
            writer = PdfWriter()
            for page_num in selected_pages:
                if 1 <= page_num <= total_pages:
                    writer.add_page(reader.pages[page_num-1])
            
            split_pdf_path = 'temp_split.pdf'
            with open(split_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(split_pdf_path, 'rb') as f:
                split_data = f.read()
            
            response = HttpResponse(
                split_data,
                content_type='application/pdf'
            )
            response['Content-Disposition'] = 'attachment; filename="split_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_split.pdf'):
                os.unlink('temp_split.pdf')
    
    return render(request, 'split_pdf.html')

def parse_page_selection(pages_input, total_pages):
    """Parse page selection like '1,3,5-8' into list of page numbers"""
    if not pages_input:
        return list(range(1, total_pages + 1))
    
    pages = []
    parts = pages_input.split(',')
    
    for part in parts:
        if '-' in part:
            start_end = part.split('-')
            if len(start_end) == 2:
                start = int(start_end[0].strip())
                end = int(start_end[1].strip())
                pages.extend(range(start, end + 1))
        else:
            pages.append(int(part.strip()))
    
    return [p for p in pages if 1 <= p <= total_pages]

# Protect PDF
def protect_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        password = request.POST.get('password', '')
        
        if not password:
            return HttpResponse("Password is required", status=400)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            writer.encrypt(password)
            
            protected_pdf_path = 'temp_protected.pdf'
            with open(protected_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(protected_pdf_path, 'rb') as f:
                protected_data = f.read()
            
            response = HttpResponse(
                protected_data,
                content_type='application/pdf'
            )
            response['Content-Disposition'] = 'attachment; filename="protected_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_protected.pdf'):
                os.unlink('temp_protected.pdf')
    
    return render(request, 'protect_pdf.html')

# Unlock PDF
def unlock_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        password = request.POST.get('password', '')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            if reader.is_encrypted:
                if not password:
                    return HttpResponse("Password is required to unlock PDF", status=400)
                reader.decrypt(password)
            
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            
            unlocked_pdf_path = 'temp_unlocked.pdf'
            with open(unlocked_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(unlocked_pdf_path, 'rb') as f:
                unlocked_data = f.read()
            
            response = HttpResponse(
                unlocked_data,
                content_type='application/pdf'
            )
            response['Content-Disposition'] = 'attachment; filename="unlocked_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_unlocked.pdf'):
                os.unlink('temp_unlocked.pdf')
    
    return render(request, 'unlock_pdf.html')

# Image to PDF
def image_to_pdf(request):
    if request.method == 'POST' and request.FILES.getlist('image_files'):
        image_files = request.FILES.getlist('image_files')
        
        if not image_files:
            return HttpResponse("Please select at least one image file", status=400)
        
        temp_files = []
        
        try:
            image_data = []
            for image_file in image_files:
                temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                for chunk in image_file.chunks():
                    temp_img.write(chunk)
                temp_img.close()
                temp_files.append(temp_img.name)
                
                with open(temp_img.name, 'rb') as f:
                    image_data.append(f.read())
            
            # Convert images to PDF
            pdf_bytes = img2pdf.convert(image_data)
            
            response = HttpResponse(
                pdf_bytes,
                content_type='application/pdf'
            )
            response['Content-Disposition'] = 'attachment; filename="converted_from_images.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            for temp_file in temp_files:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
    
    return render(request, 'image_to_pdf.html')

# Home Page
def home(request):
    return render(request, 'home.html')

# All Tools Page
def all_tools(request):
    return render(request, 'all_tools.html')


# Edit PDF - Add text and images
def edit_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        text_to_add = request.POST.get('text', '')
        position_x = int(request.POST.get('position_x', 100))
        position_y = int(request.POST.get('position_y', 100))
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Add text to PDF
            if text_to_add:
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.drawString(position_x, position_y, text_to_add)
                can.save()
                packet.seek(0)
                new_pdf = PdfReader(packet)
                writer.add_page(new_pdf.pages[0])
            
            edited_pdf_path = 'temp_edited.pdf'
            with open(edited_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(edited_pdf_path, 'rb') as f:
                edited_data = f.read()
            
            response = HttpResponse(edited_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="edited_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_edited.pdf'):
                os.unlink('temp_edited.pdf')
    
    return render(request, 'edit_pdf.html')

# Compress PDF
def compress_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                page.compress_content_streams()  # Compress the page
                writer.add_page(page)
            
            compressed_pdf_path = 'temp_compressed.pdf'
            with open(compressed_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(compressed_pdf_path, 'rb') as f:
                compressed_data = f.read()
            
            response = HttpResponse(compressed_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="compressed_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_compressed.pdf'):
                os.unlink('temp_compressed.pdf')
    
    return render(request, 'compress_pdf.html')

# Rotate PDF
def rotate_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        rotation_angle = int(request.POST.get('rotation', 90))
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                page.rotate(rotation_angle)
                writer.add_page(page)
            
            rotated_pdf_path = 'temp_rotated.pdf'
            with open(rotated_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(rotated_pdf_path, 'rb') as f:
                rotated_data = f.read()
            
            response = HttpResponse(rotated_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="rotated_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_rotated.pdf'):
                os.unlink('temp_rotated.pdf')
    
    return render(request, 'rotate_pdf.html')

# Add Watermark
def add_watermark(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        watermark_text = request.POST.get('watermark_text', 'CONFIDENTIAL')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                # Create watermark
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.setFont("Helvetica", 40)
                can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)  # Gray with transparency
                can.rotate(45)
                can.drawString(200, 100, watermark_text)
                can.save()
                packet.seek(0)
                watermark_pdf = PdfReader(packet)
                
                # Merge watermark with original page
                page.merge_page(watermark_pdf.pages[0])
                writer.add_page(page)
            
            watermarked_pdf_path = 'temp_watermarked.pdf'
            with open(watermarked_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(watermarked_pdf_path, 'rb') as f:
                watermarked_data = f.read()
            
            response = HttpResponse(watermarked_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="watermarked_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_watermarked.pdf'):
                os.unlink('temp_watermarked.pdf')
    
    return render(request, 'add_watermark.html')

# Remove Watermark (Basic implementation)
def remove_watermark(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            # This is a basic implementation - watermark removal is complex
            # For now, we'll just return the original PDF
            with open(temp_pdf_path, 'rb') as f:
                original_data = f.read()
            
            response = HttpResponse(original_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="cleaned_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
    
    return render(request, 'remove_watermark.html')

# PDF to Text
def pdf_to_text(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            text_content = ""
            with pdfplumber.open(temp_pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content += f"--- Page {page.page_number} ---\n{text}\n\n"
            
            if not text_content.strip():
                text_content = "No text could be extracted from the PDF."
            
            response = HttpResponse(text_content, content_type='text/plain')
            response['Content-Disposition'] = 'attachment; filename="extracted_text.txt"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
    
    return render(request, 'pdf_to_text.html')

# HTML to PDF
def html_to_pdf(request):
    if request.method == 'POST':
        html_content = request.POST.get('html_content', '')
        css_content = request.POST.get('css_content', '')
        
        if not html_content:
            return HttpResponse("HTML content is required", status=400)
        
        try:
            # Create PDF from HTML
            buffer = io.BytesIO()
            pisa_status = pisa.CreatePDF(html_content, dest=buffer)
            
            if pisa_status.err:
                return HttpResponse("Error generating PDF", status=500)
            
            pdf_data = buffer.getvalue()
            buffer.close()
            
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="converted_from_html.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
    
    return render(request, 'html_to_pdf.html')

# PDF Metadata Editor
def edit_metadata(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        title = request.POST.get('title', '')
        author = request.POST.get('author', '')
        subject = request.POST.get('subject', '')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Update metadata
            if title:
                writer.add_metadata({'/Title': title})
            if author:
                writer.add_metadata({'/Author': author})
            if subject:
                writer.add_metadata({'/Subject': subject})
            
            metadata_pdf_path = 'temp_metadata.pdf'
            with open(metadata_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(metadata_pdf_path, 'rb') as f:
                metadata_data = f.read()
            
            response = HttpResponse(metadata_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="metadata_updated.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_metadata.pdf'):
                os.unlink('temp_metadata.pdf')
    
    return render(request, 'edit_metadata.html')

# Page Rearrange
def rearrange_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        page_order = request.POST.get('page_order', '')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            total_pages = len(reader.pages)
            
            # Parse page order (e.g., "3,1,2" or "1-3,5,4")
            if page_order:
                pages_to_include = parse_page_selection(page_order, total_pages)
            else:
                pages_to_include = list(range(1, total_pages + 1))
            
            writer = PdfWriter()
            for page_num in pages_to_include:
                if 1 <= page_num <= total_pages:
                    writer.add_page(reader.pages[page_num-1])
            
            rearranged_pdf_path = 'temp_rearranged.pdf'
            with open(rearranged_pdf_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            with open(rearranged_pdf_path, 'rb') as f:
                rearranged_data = f.read()
            
            response = HttpResponse(rearranged_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="rearranged_document.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            if os.path.exists('temp_rearranged.pdf'):
                os.unlink('temp_rearranged.pdf')
    
    return render(request, 'rearrange_pdf.html')

# PDF Forms Fill & Edit (Basic)
def fill_pdf_form(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            # For now, return the original PDF
            # Advanced form filling requires more complex logic
            with open(temp_pdf_path, 'rb') as f:
                pdf_data = f.read()
            
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="filled_form.pdf"'
            return response
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
    
    return render(request, 'fill_pdf_form.html')

# Batch Processing (Multiple files)
def batch_process(request):
    if request.method == 'POST' and request.FILES.getlist('pdf_files'):
        pdf_files = request.FILES.getlist('pdf_files')
        operation = request.POST.get('operation', 'merge')
        
        if not pdf_files:
            return HttpResponse("Please select PDF files", status=400)
        
        temp_files = []
        
        try:
            if operation == 'merge':
                # Merge all PDFs
                merger = PdfMerger()
                for pdf_file in pdf_files:
                    temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                    for chunk in pdf_file.chunks():
                        temp_pdf.write(chunk)
                    temp_pdf.close()
                    temp_files.append(temp_pdf.name)
                    merger.append(temp_pdf.name)
                
                batch_pdf_path = 'temp_batch.pdf'
                merger.write(batch_pdf_path)
                merger.close()
                
                with open(batch_pdf_path, 'rb') as f:
                    batch_data = f.read()
                
                response = HttpResponse(batch_data, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="batch_processed.pdf"'
                return response
                
            else:
                return HttpResponse("Selected operation not yet implemented", status=400)
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            for temp_file in temp_files:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            if os.path.exists('temp_batch.pdf'):
                os.unlink('temp_batch.pdf')
    
    return render(request, 'batch_process.html')

# PDF Info Viewer
def pdf_info(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            reader = PdfReader(temp_pdf_path)
            info = reader.metadata
            
            pdf_info_data = {
                'total_pages': len(reader.pages),
                'title': info.get('/Title', 'Not specified'),
                'author': info.get('/Author', 'Not specified'),
                'subject': info.get('/Subject', 'Not specified'),
                'creator': info.get('/Creator', 'Not specified'),
                'producer': info.get('/Producer', 'Not specified'),
                'creation_date': info.get('/CreationDate', 'Not specified'),
                'modification_date': info.get('/ModDate', 'Not specified'),
                'encrypted': reader.is_encrypted,
            }
            
            return render(request, 'pdf_info.html', {'pdf_info': pdf_info_data})
            
        except Exception as e:
            return HttpResponse(f"Error: {str(e)}", status=500)
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
    
    return render(request, 'pdf_info.html')

# OCR PDF - All Formats with Layout Preservation
def ocr_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        # Import required modules
        from docx.shared import RGBColor
        from pptx.util import Inches
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        
        pdf_file = request.FILES['pdf_file']
        output_format = request.POST.get('output_format', 'searchable_pdf')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        try:
            # Convert PDF to images
            images = convert_from_path(temp_pdf_path, dpi=300)
            
            if output_format == 'searchable_pdf':
                # Searchable PDF with original layout + editable text
                output_buffer = io.BytesIO()
                pdf_canvas = canvas.Canvas(output_buffer, pagesize=A4)
                
                for page_num, image in enumerate(images):
                    if page_num > 0:
                        pdf_canvas.showPage()
                    
                    # Get OCR data with exact coordinates
                    ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    # Add text layer at exact positions
                    for i in range(len(ocr_data['text'])):
                        text = ocr_data['text'][i].strip()
                        confidence = ocr_data['conf'][i]
                        
                        if text and len(text) > 1 and confidence > 30:
                            # Get coordinates
                            x = ocr_data['left'][i]
                            y = image.height - ocr_data['top'][i] - ocr_data['height'][i]
                            
                            # Convert coordinates to PDF scale
                            pdf_x = x * (A4[0] / image.width)
                            pdf_y = y * (A4[1] / image.height)
                            
                            # Calculate font size
                            font_size = max(8, ocr_data['height'][i] * (A4[1] / image.height) * 0.7)
                            
                            # Add invisible but selectable text
                            pdf_canvas.setFont("Helvetica", font_size)
                            pdf_canvas.setFillColorRGB(0, 0, 0, alpha=0)  # Completely transparent
                            pdf_canvas.drawString(pdf_x, pdf_y, text)
                
                pdf_canvas.save()
                output_buffer.seek(0)
                
                # Merge with original PDF
                original_pdf = PdfReader(temp_pdf_path)
                text_layer_pdf = PdfReader(output_buffer)
                
                writer = PdfWriter()
                
                for page_num in range(len(original_pdf.pages)):
                    original_page = original_pdf.pages[page_num]
                    if page_num < len(text_layer_pdf.pages):
                        text_page = text_layer_pdf.pages[page_num]
                        original_page.merge_page(text_page)
                    
                    writer.add_page(original_page)
                
                # Save final PDF
                final_output_path = 'temp_searchable.pdf'
                with open(final_output_path, 'wb') as f:
                    writer.write(f)
                
                with open(final_output_path, 'rb') as f:
                    output_data = f.read()
                
                response = HttpResponse(output_data, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="searchable_document.pdf"'
                return response
                
            elif output_format == 'word':
                # Word with embedded images + text
                from docx import Document
                doc = Document()
                
                for page_num, image in enumerate(images):
                    # Add page break for subsequent pages
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # Convert image to bytes for Word
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='PNG')
                    img_byte_arr.seek(0)
                    
                    # Add image to Word (original layout)
                    doc.add_picture(img_byte_arr, width=Inches(6.5))
                    
                    # Get OCR text
                    ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    # Add paragraph with OCR text (white text)
                    text_paragraph = ""
                    for i in range(len(ocr_data['text'])):
                        text = ocr_data['text'][i].strip()
                        confidence = ocr_data['conf'][i]
                        if text and confidence > 30:
                            text_paragraph += text + " "
                    
                    # Add hidden text paragraph
                    if text_paragraph.strip():
                        p = doc.add_paragraph()
                        run = p.add_run(text_paragraph.strip())
                        run.font.color.rgb = RGBColor(255, 255, 255)  # White color
                
                temp_output_path = 'temp_ocr_output.docx'
                doc.save(temp_output_path)
                
                with open(temp_output_path, 'rb') as f:
                    output_data = f.read()
                
                response = HttpResponse(output_data, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                response['Content-Disposition'] = 'attachment; filename="ocr_document.docx"'
                return response
                
            elif output_format == 'excel':
                # Excel with detailed OCR data
                workbook = Workbook()
                
                for page_num, image in enumerate(images):
                    # Create sheet for each page
                    if page_num == 0:
                        worksheet = workbook.active
                        worksheet.title = f"Page_{page_num+1}"
                    else:
                        worksheet = workbook.create_sheet(title=f"Page_{page_num+1}")
                    
                    # Add headers
                    worksheet.cell(row=1, column=1, value="Text")
                    worksheet.cell(row=1, column=2, value="Confidence")
                    worksheet.cell(row=1, column=3, value="Position")
                    
                    # Get detailed OCR data
                    ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    row_num = 2
                    for i in range(len(ocr_data['text'])):
                        text = ocr_data['text'][i].strip()
                        confidence = ocr_data['conf'][i]
                        
                        if text and confidence > 30:
                            worksheet.cell(row=row_num, column=1, value=text)
                            worksheet.cell(row=row_num, column=2, value=confidence)
                            worksheet.cell(row=row_num, column=3, value=f"({ocr_data['left'][i]},{ocr_data['top'][i]})")
                            row_num += 1
                
                temp_output_path = 'temp_ocr_output.xlsx'
                workbook.save(temp_output_path)
                
                with open(temp_output_path, 'rb') as f:
                    output_data = f.read()
                
                response = HttpResponse(output_data, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                response['Content-Disposition'] = 'attachment; filename="ocr_data.xlsx"'
                return response
                
            elif output_format == 'ppt':
                # PowerPoint with images + OCR text - FIXED VERSION
                from pptx import Presentation
                from pptx.dml.color import RGBColor as PptRGBColor
                
                prs = Presentation()
                
                for page_num, image in enumerate(images):
                    # Create slide with blank layout
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Convert image to bytes
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='PNG')
                    img_byte_arr.seek(0)
                    
                    # Add image to slide (original layout)
                    left = Inches(0.5)
                    top = Inches(0.5)
                    pic = slide.shapes.add_picture(img_byte_arr, left, top, height=Inches(7))
                    
                    # Get OCR text
                    ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    # Add text box with OCR text
                    text_paragraph = ""
                    for i in range(len(ocr_data['text'])):
                        text = ocr_data['text'][i].strip()
                        confidence = ocr_data['conf'][i]
                        if text and confidence > 30:
                            text_paragraph += text + " "
                    
                    if text_paragraph.strip():
                        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(9), Inches(1))
                        text_frame = textbox.text_frame
                        text_frame.text = text_paragraph.strip()
                        
                        # FIX: PowerPoint ke liye proper white text
                        for paragraph in text_frame.paragraphs:
                            # Clear any existing runs
                            paragraph.clear()
                            # Add new run with white text
                            run = paragraph.add_run()
                            run.text = text_paragraph.strip()
                            run.font.color.rgb = PptRGBColor(255, 255, 255)  # PowerPoint RGBColor
                
                temp_output_path = 'temp_ocr_output.pptx'
                prs.save(temp_output_path)
                
                with open(temp_output_path, 'rb') as f:
                    output_data = f.read()
                
                response = HttpResponse(output_data, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                response['Content-Disposition'] = 'attachment; filename="ocr_presentation.pptx"'
                return response
            
        except Exception as e:
            return HttpResponse(f"OCR Error: {str(e)}", status=500)
        finally:
            # Cleanup
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)
            for temp_file in ['temp_searchable.pdf', 'temp_ocr_output.docx', 'temp_ocr_output.xlsx', 'temp_ocr_output.pptx']:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
    
    return render(request, 'ocr_pdf.html')

def about(request):
    return render(request, 'about.html')

def contact(request):
    return render(request, 'contact.html')

def careers(request):
    return render(request, 'careers.html')

def blog(request):
    return render(request, 'blog.html')

def affiliate(request):
    return render(request, 'affiliate.html')

def privacy_policy(request):
    return render(request, 'privacy_policy.html')

def terms_of_service(request):
    return render(request, 'terms_of_service.html')

def disclaimer(request):
    return render(request, 'disclaimer.html')

def cookie_policy(request):
    return render(request, 'cookie_policy.html')

def gdpr(request):
    return render(request, 'gdpr.html')